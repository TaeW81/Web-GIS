from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "1km"
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = RGBColor(0, 0, 0)

# Try adding glow
rPr = run._r.get_or_add_rPr()
glow_xml = """
<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
    <a:glow rad="50800">
        <a:srgbClr val="FFFFFF"/>
    </a:glow>
</a:effectLst>
"""
effectLst = parse_xml(glow_xml)
rPr.append(effectLst)

# Try adding outline
outline_xml = """
<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" w="12700">
    <a:solidFill>
        <a:srgbClr val="FFFFFF"/>
    </a:solidFill>
</a:ln>
"""
ln = parse_xml(outline_xml)
rPr.append(ln)

prs.save("test_text_outline.pptx")
print("Saved")
