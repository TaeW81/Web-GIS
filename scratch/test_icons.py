from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# North Arrow
arrow = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(8), Inches(1), Inches(0.4), Inches(0.6))
arrow.fill.solid()
arrow.fill.fore_color.rgb = RGBColor(0, 0, 0)
arrow.line.color.rgb = RGBColor(255, 255, 255)
arrow.line.width = Pt(1.5)

n_text = slide.shapes.add_textbox(Inches(7.9), Inches(0.5), Inches(0.6), Inches(0.4))
n_text.text_frame.text = "N"
n_text.text_frame.paragraphs[0].font.bold = True
n_text.text_frame.paragraphs[0].font.size = Pt(20)
n_text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Highway Shield
hw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(0.8), Inches(0.5))
hw.fill.solid()
hw.fill.fore_color.rgb = RGBColor(0, 51, 160) # Blue
hw.line.color.rgb = RGBColor(255, 255, 255)
hw.text_frame.text = "고속"
hw.text_frame.paragraphs[0].font.size = Pt(10)
hw.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
hw.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# National Road
nr = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(1), Inches(0.6), Inches(0.5))
nr.fill.solid()
nr.fill.fore_color.rgb = RGBColor(255, 204, 0) # Yellow
nr.line.color.rgb = RGBColor(0, 51, 160)
nr.text_frame.text = "국도"
nr.text_frame.paragraphs[0].font.size = Pt(10)
nr.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
nr.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Mountain
mt = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3), Inches(1), Inches(0.5), Inches(0.5))
mt.fill.solid()
mt.fill.fore_color.rgb = RGBColor(34, 139, 34) # Green
mt.line.color.rgb = RGBColor(255, 255, 255)

# Place names
tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(1), Inches(0.5))
tb.text_frame.text = "OO동"
tb.text_frame.paragraphs[0].font.size = Pt(14)
tb.text_frame.paragraphs[0].font.bold = True

prs.save("test_icons.pptx")
print("Saved")
