import io
import math
import requests
from PIL import Image
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from shapely.geometry import Polygon
import sys
import os

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import VWORLD_KEY, MAP_SOURCES, VWORLD_DOMAIN

def deg2num(lat_deg, lon_deg, zoom):
    lat_rad = math.radians(lat_deg)
    n = 2.0 ** zoom
    xtile = int((lon_deg + 180.0) / 360.0 * n)
    ytile = int((1.0 - math.asinh(math.tan(lat_rad)) / math.pi) / 2.0 * n)
    return (xtile, ytile)

def num2deg(xtile, ytile, zoom):
    n = 2.0 ** zoom
    lon_deg = xtile / n * 360.0 - 180.0
    lat_rad = math.atan(math.sinh(math.pi * (1 - 2 * ytile / n)))
    lat_deg = math.degrees(lat_rad)
    return (lat_deg, lon_deg)

def create_qbs_pptx(boundary_polygon, visible_layers):
    min_x, min_y, max_x, max_y = boundary_polygon.bounds
    center_x = (min_x + max_x) / 2
    center_y = (min_y + max_y) / 2
    
    # 1. Tile logic
    zoom = 13
    radius_km = 6.0
    lat_km = 111.32
    lon_km = 111.32 * math.cos(math.radians(center_y))
    dx = radius_km / lon_km
    dy = radius_km / lat_km
    
    exp_min_x = center_x - dx
    exp_max_x = center_x + dx
    exp_min_y = center_y - dy
    exp_max_y = center_y + dy
    
    x_min_tile, y_max_tile = deg2num(exp_min_y, exp_min_x, zoom) # min_y -> max_ytile
    x_max_tile, y_min_tile = deg2num(exp_max_y, exp_max_x, zoom) # max_y -> min_ytile
    
    num_x = x_max_tile - x_min_tile + 1
    num_y = y_max_tile - y_min_tile + 1
    
    bg_img = Image.new("RGBA", (num_x * 256, num_y * 256), (255, 255, 255, 255))
    
    for i, x in enumerate(range(x_min_tile, x_max_tile + 1)):
        for j, y in enumerate(range(y_min_tile, y_max_tile + 1)):
            url = f"https://api.vworld.kr/req/wmts/1.0.0/{VWORLD_KEY}/Base/{zoom}/{y}/{x}.png"
            try:
                res = requests.get(url, timeout=5)
                if res.status_code == 200:
                    tile = Image.open(io.BytesIO(res.content)).convert("RGBA")
                    bg_img.paste(tile, (i * 256, j * 256))
            except Exception as e:
                pass
                
    # Calculate exact bbox of the stitched image
    img_max_lat, img_min_lon = num2deg(x_min_tile, y_min_tile, zoom)
    img_min_lat, img_max_lon = num2deg(x_max_tile + 1, y_max_tile + 1, zoom)
    
    # 2. WMS Overlays
    wms_url = "http://api.vworld.kr/req/wms"
    base_params = {
        "key": VWORLD_KEY, "domain": VWORLD_DOMAIN,
        "service": "WMS", "request": "GetMap",
        "crs": "EPSG:4326", "format": "image/png", 
        "width": str(num_x * 256), "height": str(num_y * 256),
        "bbox": f"{img_min_lon},{img_min_lat},{img_max_lon},{img_max_lat}",
        "transparent": "true"
    }
    
    layer_codes = []
    for source_name, categories in MAP_SOURCES.items():
        if source_name != "브이월드 (VWorld)":
            continue
        for cat_name, layers in categories.items():
            for name, code in layers.items():
                if name in visible_layers and not "READY" in str(code):
                    layer_codes.append(code.lower())
    
    for code in layer_codes:
        try:
            ol_params = base_params.copy()
            ol_params["layers"] = code
            res_ol = requests.get(wms_url, params=ol_params, timeout=10)
            if res_ol.status_code == 200:
                ol_img = Image.open(io.BytesIO(res_ol.content)).convert("RGBA")
                bg_img = Image.alpha_composite(bg_img, ol_img)
        except Exception as e:
            print(f"Failed to overlay {code}: {e}")
            
    # 3. Plotting using Matplotlib
    fig, ax = plt.subplots(figsize=(10, 10))
    ax.imshow(bg_img, extent=[img_min_lon, img_max_lon, img_min_lat, img_max_lat], origin='upper')

    # Draw boundary
    for sub_poly in ([boundary_polygon] if boundary_polygon.geom_type == 'Polygon' else boundary_polygon.geoms):
        bx, by = sub_poly.exterior.xy
        ax.plot(bx, by, color='red', lw=2.0, zorder=10)
        ax.fill(bx, by, color='red', alpha=0.3, zorder=9)
        
    ax.set_xlim(exp_min_x, exp_max_x)
    ax.set_ylim(exp_min_y, exp_max_y)
    ax.set_aspect('equal')
    ax.axis('off')
    
    map_buf = io.BytesIO()
    plt.savefig(map_buf, format='png', dpi=200, bbox_inches='tight', pad_inches=0)
    plt.close(fig)
    map_buf.seek(0)
    
    # 4. PPTX Creation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    map_size = Inches(7.5)
    map_left = (prs.slide_width - map_size) / 2
    map_top = 0
    slide.shapes.add_picture(map_buf, map_left, map_top, width=map_size, height=map_size)
    
    center_px = map_left + map_size / 2
    center_py = map_top + map_size / 2
    
    for r_km in [1, 2, 3, 4, 5]:
        r_inch = map_size * (r_km / (radius_km * 2))
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, 
            center_px - r_inch, 
            center_py - r_inch, 
            r_inch * 2, 
            r_inch * 2
        )
        shape.fill.background()
        shape.line.color.rgb = RGBColor(128, 128, 128)
        if r_km % 2 == 1:
            shape.line.dash_style = 4
        
        txBox = slide.shapes.add_textbox(
            center_px - r_inch - Inches(0.4), 
            center_py - Inches(0.15), 
            Inches(0.8), 
            Inches(0.3)
        )
        tf = txBox.text_frame
        tf.text = f"{r_km}km"
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = RGBColor(255, 255, 255)
        txBox.line.color.rgb = RGBColor(128, 128, 128)

    target_box = slide.shapes.add_textbox(
        center_px - Inches(0.6), 
        center_py - Inches(0.4), 
        Inches(1.2), 
        Inches(0.4)
    )
    target_box.fill.solid()
    target_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    target_box.line.color.rgb = RGBColor(0, 0, 0)
    tf = target_box.text_frame
    tf.text = "사업대상지"
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    prs.save("test_qbs.pptx")
    print("test_qbs.pptx created successfully")

if __name__ == "__main__":
    lon, lat = 126.9780, 37.5665
    coords = [(lon-0.005, lat-0.005), (lon+0.005, lat-0.005), (lon+0.005, lat+0.005), (lon-0.005, lat+0.005)]
    poly = Polygon(coords)
    create_qbs_pptx(poly, ["지적도", "토지이용계획도"])
