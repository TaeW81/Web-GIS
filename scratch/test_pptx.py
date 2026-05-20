import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank slide

# Add a circle
cx, cy = Inches(5), Inches(3.75)
radius = Inches(1)
left = cx - radius
top = cy - radius
width = height = radius * 2

shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
shape.fill.background()
shape.line.color.rgb = RGBColor(255, 0, 0)

prs.save("test.pptx")
print("Saved test.pptx")
