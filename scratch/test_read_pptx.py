from pptx import Presentation
import os

filepath = os.path.join("자료", "qbs_위치도.pptx")
prs = Presentation(filepath)
print("Slides:", len(prs.slides))
for i, slide in enumerate(prs.slides):
    print(f"Slide {i} shapes: {len(slide.shapes)}")
    for j, shape in enumerate(slide.shapes):
        print(f"  Shape {j}: type={shape.shape_type}, name={shape.name}")
