import copy
from pptx import Presentation

prs = Presentation('template.pptx')
source_slide = prs.slides[1]
target_slide = prs.slides[0]

# shape 0 is Group 1 (Highway)
group_shape = source_slide.shapes[0]

print(group_shape.shape_type)
el = copy.deepcopy(group_shape.element)
# change position
el.grpSpPr.xfrm.off.x = 1000000
el.grpSpPr.xfrm.off.y = 1000000

# Try to find text element to change road number
for shape_el in el.sp_lst:
    # Look for txBody
    txBody = shape_el.find('.//a:txBody', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
    if txBody is not None:
        for t in txBody.findall('.//a:t', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
            if t.text.strip().isdigit():
                print("Found text:", t.text)
                t.text = "99"

target_slide.shapes._spTree.append(el)
prs.save("test_copy.pptx")
print("Saved")
