"""
QBS 위치도 (Site Location Map) 생성기 — 보고서 품질 버전

설계 방침:
  1. 모든 지리 레이어(배경/도로/철도/하천/사업지)를 matplotlib 단일 figure에서 통합 렌더링
     → 픽셀 단위 정밀 제어 (선 굵기/색/투명도/zorder)
  2. 배경 위성지도는 그레이스케일 + 밝기 보정 → 전경(도로/철도/사업지)이 잘 보이도록
  3. 도로 등급별 색상/굵기 차별화 (고속도로 → 국도 → 지방도)
  4. 철도는 점선 + 진한 색 (배경과 명확히 구분)
  5. 사업지: 빨간 음영 폴리곤 + 굵은 테두리 → 한눈에 인식
  6. 5km/10km 점선 동심원 (점선 회색)
  7. 결과 고해상도 PNG → PPT 슬라이드에 메인 이미지로 삽입
  8. PPT 위에 추가 도형: 사업대상지 라벨박스, 방위표, 스케일바
     (사용자가 PPT에서 위치/텍스트 수정 가능)

데이터 소스:
  - 배경 타일: V-World WMTS Base
  - 도로: V-World WFS lt_l_highway (고속도로) / lt_l_natlroad (국도) / lt_l_moctlink (교통링크)
  - 철도: V-World WFS lt_l_frstrail
  - 하천: V-World WFS lt_c_wkmstrm
  - POI: V-World 검색 API (IC/JC/역)

향후 개선 후보:
  - OSM Overpass API 추가 (노선별 색상 정확화, 지하철 호선 표시)
  - SVG 벡터 export (Illustrator 편집)
"""
import io
import math
import requests
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Circle, FancyBboxPatch
from matplotlib.lines import Line2D
from PIL import Image, ImageEnhance, ImageOps
from pyproj import Transformer
from pptx import Presentation
from pptx.util import Cm, Pt, Emu, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from config import VWORLD_KEY, VWORLD_DOMAIN


# ──────────────────────────────────────────────────────────────────
# 타일 좌표 변환
# ──────────────────────────────────────────────────────────────────
def deg2num(lat_deg, lon_deg, zoom):
    lat_rad = math.radians(lat_deg)
    n = 2.0 ** zoom
    return (int((lon_deg + 180.0) / 360.0 * n),
            int((1.0 - math.asinh(math.tan(lat_rad)) / math.pi) / 2.0 * n))


def num2deg(xtile, ytile, zoom):
    n = 2.0 ** zoom
    lon_deg = xtile / n * 360.0 - 180.0
    lat_rad = math.atan(math.sinh(math.pi * (1 - 2 * ytile / n)))
    return (math.degrees(lat_rad), lon_deg)


class QBSGenerator:
    """사업지 위치도 (Site Location Map) PPT 생성기."""

    # === 시각 디자인 토큰 ============================================
    BG_DESATURATION = 0.4       # 배경 채도 (낮을수록 회색)
    BG_BRIGHTNESS = 1.15        # 배경 밝기 (>1.0 = 더 밝게)

    # 도로 색상/굵기 (등급별)
    HW_OUTLINE_COLOR = "#ffffff"
    HW_FILL_COLOR = "#2c3e50"   # 고속도로 (진한 남색)
    HW_OUTLINE_LW = 5.0
    HW_FILL_LW = 3.0

    NATL_OUTLINE_COLOR = "#ffffff"
    NATL_FILL_COLOR = "#444444"  # 국도 (진회색)
    NATL_OUTLINE_LW = 4.0
    NATL_FILL_LW = 2.2

    LOCAL_COLOR = "#999999"      # 지방도/일반도로 (옅은 회색)
    LOCAL_LW = 1.0

    # 철도 (점선)
    RAIL_OUTLINE_COLOR = "#ffffff"
    RAIL_FILL_COLOR = "#1a472a"  # 진한 녹색
    RAIL_OUTLINE_LW = 5.0
    RAIL_FILL_LW = 2.5

    # 하천
    WATER_FILL = "#bbdefb"
    WATER_EDGE = "#64b5f6"

    # 사업지
    SITE_FILL = "#ff0000"
    SITE_EDGE = "#b71c1c"
    SITE_ALPHA = 0.45
    SITE_LW = 2.5
    # =============================================================

    def __init__(self, boundary_polygon, visible_layers=None):
        self.boundary_polygon = boundary_polygon
        self.visible_layers = visible_layers or []

    # ─────────────────────────────────────────────────────────────
    # 배경 타일 다운로드 + 그레이스케일 + 밝기 보정
    # ─────────────────────────────────────────────────────────────
    def _fetch_base(self, x0, y0, x1, y1, zoom):
        nx, ny = x1 - x0 + 1, y1 - y0 + 1
        base = Image.new("RGBA", (nx * 256, ny * 256), (255, 255, 255, 255))
        for i, x in enumerate(range(x0, x1 + 1)):
            for j, y in enumerate(range(y0, y1 + 1)):
                try:
                    url = f"http://api.vworld.kr/req/wmts/1.0.0/{VWORLD_KEY}/Base/{zoom}/{y}/{x}.png"
                    r = requests.get(url, timeout=5)
                    if r.status_code == 200:
                        tile = Image.open(io.BytesIO(r.content)).convert("RGBA")
                        base.paste(tile, (i * 256, j * 256))
                except Exception:
                    pass

        # 그레이스케일 변환 (채도 낮추기)
        rgb = base.convert("RGB")
        gray = ImageOps.grayscale(rgb).convert("RGB")
        blended = Image.blend(rgb, gray, 1.0 - self.BG_DESATURATION)
        # 밝기 살짝 ↑ (전경 가독성)
        blended = ImageEnhance.Brightness(blended).enhance(self.BG_BRIGHTNESS)
        # 대비도 약간 ↓
        blended = ImageEnhance.Contrast(blended).enhance(0.85)
        return blended.convert("RGBA")

    # ─────────────────────────────────────────────────────────────
    # V-World WFS 호출 헬퍼
    # ─────────────────────────────────────────────────────────────
    @staticmethod
    def _wfs(typename, bbox3857, maxf=2000):
        try:
            return requests.get(
                "http://api.vworld.kr/req/wfs",
                params={
                    "key": VWORLD_KEY, "SERVICE": "WFS", "version": "1.1.0",
                    "request": "GetFeature", "TYPENAME": typename,
                    "BBOX": bbox3857, "outputFormat": "application/json",
                    "maxFeatures": maxf,
                },
                timeout=20,
            )
        except Exception:
            return None

    @staticmethod
    def _lines(feature):
        g = feature["geometry"]
        cs = g["coordinates"]
        return cs if g["type"] == "MultiLineString" else [cs]

    # ─────────────────────────────────────────────────────────────
    # POI 검색 (IC/JC/역)
    # ─────────────────────────────────────────────────────────────
    @staticmethod
    def _search_pois(bbox_str, query, size=50):
        try:
            url = (
                f"http://api.vworld.kr/req/search?service=search&request=search"
                f"&version=2.0&crs=EPSG:4326&bbox={bbox_str}&type=PLACE"
                f"&query={query}&key={VWORLD_KEY}&size={size}&format=json"
            )
            r = requests.get(url, timeout=10)
            if r.status_code == 200:
                data = r.json()
                if "result" in data.get("response", {}):
                    return data["response"]["result"]["items"]
        except Exception:
            pass
        return []

    # ─────────────────────────────────────────────────────────────
    # 메인 생성 함수
    # ─────────────────────────────────────────────────────────────
    def generate(self):
        # ── 1. 영역 계산 ──
        min_x, min_y, max_x, max_y = self.boundary_polygon.bounds
        cx, cy = (min_x + max_x) / 2, (min_y + max_y) / 2

        # 슬라이드 크기 (cm) — 세로형 위치도
        SW_CM, SH_CM = 16.0, 14.0
        SW, SH = SW_CM / 2.54, SH_CM / 2.54  # inch

        # 표시 반경 — 10km 원이 들어오도록 12km
        R_KM_Y = 12.0
        R_KM_X = R_KM_Y * (SW_CM / SH_CM)
        lat_km = 111.32
        lon_km = 111.32 * math.cos(math.radians(cy))
        dx_deg = R_KM_X / lon_km
        dy_deg = R_KM_Y / lat_km
        eMinX = cx - dx_deg; eMaxX = cx + dx_deg
        eMinY = cy - dy_deg; eMaxY = cy + dy_deg

        # ── 2. 배경 타일 ──
        zoom = 12
        x0, y0 = deg2num(eMaxY, eMinX, zoom)
        x1, y1 = deg2num(eMinY, eMaxX, zoom)
        x0 -= 1; y0 -= 1; x1 += 1; y1 += 1
        base_img = self._fetch_base(x0, y0, x1, y1, zoom)
        iMaxLat, iMinLon = num2deg(x0, y0, zoom)
        iMinLat, iMaxLon = num2deg(x1 + 1, y1 + 1, zoom)
        ext = [iMinLon, iMaxLon, iMinLat, iMaxLat]

        # ── 3. matplotlib 단일 figure에 모든 레이어 통합 ──
        plt.rcParams["font.family"] = "Malgun Gothic"
        plt.rcParams["axes.unicode_minus"] = False

        fig, ax = plt.subplots(figsize=(SW, SH), dpi=300, constrained_layout=False)
        fig.patch.set_alpha(0.0)
        ax.set_position([0, 0, 1, 1])
        ax.set_xlim(eMinX, eMaxX); ax.set_ylim(eMinY, eMaxY)
        ax.set_aspect("auto")
        ax.axis("off")

        # 3-1. 배경 (그레이스케일)
        ax.imshow(base_img, extent=ext, origin="upper", aspect="auto", zorder=1)

        # 3-2. 좌표 변환기 (EPSG:4326 ↔ 3857)
        tp = Transformer.from_crs("EPSG:4326", "EPSG:3857", always_xy=True)
        ti = Transformer.from_crs("EPSG:3857", "EPSG:4326", always_xy=True)
        m0x, m0y = tp.transform(eMinX, eMinY)
        m1x, m1y = tp.transform(eMaxX, eMaxY)
        bb3857 = f"{m0x},{m0y},{m1x},{m1y}"

        # 3-3. 하천망 (zorder 2)
        try:
            r = self._wfs("lt_c_wkmstrm", bb3857, 1500)
            if r and r.status_code == 200:
                for f in r.json().get("features", []):
                    g = f["geometry"]
                    polys = g["coordinates"] if g["type"] == "MultiPolygon" else [g["coordinates"]]
                    for poly in polys:
                        ring = poly[0] if isinstance(poly[0][0], list) else poly
                        lons, lats = zip(*[ti.transform(p[0], p[1]) for p in ring])
                        ax.fill(lons, lats, color=self.WATER_FILL, alpha=0.6, zorder=2)
                        ax.plot(lons, lats, color=self.WATER_EDGE, lw=0.8, alpha=0.8, zorder=2.1)
        except Exception:
            pass

        # 3-4. 일반 도로 (지방도) — 가장 아래 zorder=3
        try:
            r = self._wfs("lt_l_moctlink", bb3857, 5000)
            if r and r.status_code == 200:
                for f in r.json().get("features", []):
                    for ln in self._lines(f):
                        lons, lats = zip(*[ti.transform(p[0], p[1]) for p in ln])
                        ax.plot(lons, lats, color=self.LOCAL_COLOR,
                                lw=self.LOCAL_LW, alpha=0.85, zorder=3,
                                solid_capstyle="round")
        except Exception:
            pass

        # 3-5. 국도 (zorder 4)
        try:
            r = self._wfs("lt_l_natlroad", bb3857, 2000)
            if r and r.status_code == 200:
                for f in r.json().get("features", []):
                    for ln in self._lines(f):
                        lons, lats = zip(*[ti.transform(p[0], p[1]) for p in ln])
                        # 흰 외곽선 → 진회색 본선 (도로감)
                        ax.plot(lons, lats, color=self.NATL_OUTLINE_COLOR,
                                lw=self.NATL_OUTLINE_LW, alpha=1.0, zorder=4,
                                solid_capstyle="round")
                        ax.plot(lons, lats, color=self.NATL_FILL_COLOR,
                                lw=self.NATL_FILL_LW, alpha=1.0, zorder=4.1,
                                solid_capstyle="round")
        except Exception:
            pass

        # 3-6. 고속도로 (zorder 5 — 도로 중 가장 위)
        try:
            r = self._wfs("lt_l_highway", bb3857, 800)
            if r and r.status_code == 200:
                for f in r.json().get("features", []):
                    for ln in self._lines(f):
                        lons, lats = zip(*[ti.transform(p[0], p[1]) for p in ln])
                        ax.plot(lons, lats, color=self.HW_OUTLINE_COLOR,
                                lw=self.HW_OUTLINE_LW, alpha=1.0, zorder=5,
                                solid_capstyle="round")
                        ax.plot(lons, lats, color=self.HW_FILL_COLOR,
                                lw=self.HW_FILL_LW, alpha=1.0, zorder=5.1,
                                solid_capstyle="round")
        except Exception:
            pass

        # 3-7. 철도 (zorder 6 — 점선)
        try:
            r = self._wfs("lt_l_frstrail", bb3857, 1500)
            if r and r.status_code == 200:
                for f in r.json().get("features", []):
                    for ln in self._lines(f):
                        lons, lats = zip(*[ti.transform(p[0], p[1]) for p in ln])
                        # 흰 외곽 → 진녹색 점선
                        ax.plot(lons, lats, color=self.RAIL_OUTLINE_COLOR,
                                lw=self.RAIL_OUTLINE_LW, alpha=0.9, zorder=6,
                                solid_capstyle="round")
                        ax.plot(lons, lats, color=self.RAIL_FILL_COLOR,
                                lw=self.RAIL_FILL_LW, alpha=0.95, zorder=6.1,
                                linestyle=(0, (8, 4)),  # 명확한 점선
                                solid_capstyle="round")
        except Exception:
            pass

        # 3-8. 사업대상지 폴리곤 (zorder 7 — 최상단)
        polys = [self.boundary_polygon] if self.boundary_polygon.geom_type == "Polygon" \
            else list(self.boundary_polygon.geoms)
        for sp in polys:
            bx, by = sp.exterior.xy
            ax.fill(bx, by, color=self.SITE_FILL, alpha=self.SITE_ALPHA, zorder=7)
            ax.plot(bx, by, color=self.SITE_EDGE, lw=self.SITE_LW, zorder=7.1)

        # 3-9. 반경 5km / 10km 점선 동심원
        for r_km in [5, 10]:
            # lon/lat 비율 보정한 타원이 아닌 평면 거리상의 원
            # matplotlib는 axes coordinates 기준이라 ellipse 효과 발생
            # 단순화: 위도 기준 변환 사용
            theta = [i * 2 * math.pi / 180 for i in range(180 + 1)]
            xs = [cx + (r_km / lon_km) * math.cos(t) for t in theta]
            ys = [cy + (r_km / lat_km) * math.sin(t) for t in theta]
            ax.plot(xs, ys, color="#666666", lw=1.2, alpha=0.7,
                    linestyle=(0, (4, 3)), zorder=8)

        # ── 4. 결과를 PNG로 export ──
        buf_map = io.BytesIO()
        plt.savefig(buf_map, format="png", dpi=300, bbox_inches=None,
                    pad_inches=0, transparent=False, facecolor="white")
        plt.close(fig)
        buf_map.seek(0)

        # ── 5. PPT 생성 ──
        prs = Presentation()
        prs.slide_width = Cm(SW_CM); prs.slide_height = Cm(SH_CM)
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
        slide.shapes.add_picture(buf_map, 0, 0,
                                 width=Cm(SW_CM), height=Cm(SH_CM))

        # ── 6. PPT 위에 추가 도형 (편집 가능) ──
        cpx = Cm(SW_CM) / 2
        cpy = Cm(SH_CM) / 2

        # 6-1. 반경 라벨 (5km, 10km 텍스트)
        for r_km in [5, 10]:
            rpx_y = Cm(SH_CM) * (r_km / (R_KM_Y * 2))
            tb = slide.shapes.add_textbox(
                cpx + rpx_y - Cm(0.8), cpy - Cm(0.3),
                Cm(1.6), Cm(0.6),
            )
            self._zero_margins(tb)
            run = tb.text_frame.paragraphs[0].add_run()
            run.text = f"{r_km}km"
            run.font.size = Pt(10); run.font.bold = True
            run.font.color.rgb = RGBColor(50, 50, 50); run.font.name = "맑은 고딕"
            tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._text_outline_white(run)

        # 6-2. 사업대상지 라벨박스 (빨간색)
        label_w, label_h = Cm(2.8), Cm(0.75)
        lb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            cpx - label_w / 2, cpy + Cm(0.4),
            label_w, label_h,
        )
        lb.fill.solid(); lb.fill.fore_color.rgb = RGBColor(220, 30, 30)
        lb.line.color.rgb = RGBColor(255, 255, 255); lb.line.width = Pt(1.2)
        self._zero_margins(lb)
        run_lb = lb.text_frame.paragraphs[0].add_run()
        run_lb.text = "사업대상지"
        run_lb.font.size = Pt(12); run_lb.font.bold = True
        run_lb.font.color.rgb = RGBColor(255, 255, 255); run_lb.font.name = "맑은 고딕"
        lb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 6-3. 방위표 (우상단)
        self._add_compass(slide, SW_CM)

        # 6-4. 스케일바 (좌하단)
        self._add_scalebar(slide, SW_CM, SH_CM, R_KM_Y)

        # 6-5. POI 라벨 — IC/JC + 주요 역
        bbox_str = f"{eMinX},{eMinY},{eMaxX},{eMaxY}"
        drawn = set()
        self._add_pois(slide, "인터체인지", drawn, bbox_str,
                       eMinX, eMaxX, eMinY, eMaxY, SW, SH,
                       icon="ic_jc", color=(255, 140, 0))
        self._add_pois(slide, "JC", drawn, bbox_str,
                       eMinX, eMaxX, eMinY, eMaxY, SW, SH,
                       icon="ic_jc", color=(255, 140, 0))
        self._add_pois(slide, "역", drawn, bbox_str,
                       eMinX, eMaxX, eMinY, eMaxY, SW, SH,
                       icon="station", color=(40, 40, 40))

        out = io.BytesIO(); prs.save(out); out.seek(0)
        return out.getvalue()

    # ─────────────────────────────────────────────────────────────
    # PPT 헬퍼들
    # ─────────────────────────────────────────────────────────────
    @staticmethod
    def _zero_margins(shape):
        tf = shape.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

    @staticmethod
    def _text_outline_white(run):
        """텍스트에 흰색 테두리 + 글로우 (가독성)."""
        xml_ln = (
            '<a:ln w="20000" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
            '<a:round/></a:ln>'
        )
        run._r.get_or_add_rPr().append(parse_xml(xml_ln))
        run._r.get_or_add_rPr().append(parse_xml(
            '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:glow rad="25000"><a:srgbClr val="FFFFFF"/></a:glow></a:effectLst>'
        ))

    def _add_compass(self, slide, sw_cm):
        """우상단 방위표 (N)."""
        size = Cm(1.6)
        x = Cm(sw_cm) - size - Cm(0.4)
        y = Cm(0.4)

        # 흰 원 배경
        bg = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.line.color.rgb = RGBColor(80, 80, 80); bg.line.width = Pt(1.2)

        # 북쪽 화살표 (위쪽 빨간 삼각)
        aw = Cm(0.5); ah = Cm(0.7)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            x + size / 2 - aw / 2, y + Cm(0.15),
            aw, ah,
        )
        arrow.fill.solid(); arrow.fill.fore_color.rgb = RGBColor(220, 30, 30)
        arrow.line.fill.background()

        # 남쪽 회색 삼각 (역삼각)
        s_arrow_xml = (
            '<a:prstGeom prst="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<a:avLst/></a:prstGeom>'
        )
        s_arrow = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            x + size / 2 - aw / 2, y + size - ah - Cm(0.15),
            aw, ah,
        )
        s_arrow.fill.solid(); s_arrow.fill.fore_color.rgb = RGBColor(160, 160, 160)
        s_arrow.line.fill.background()
        s_arrow.rotation = 180

        # N 글자
        ntx = slide.shapes.add_textbox(x, y + Cm(0.4), size, Cm(0.5))
        self._zero_margins(ntx)
        rn = ntx.text_frame.paragraphs[0].add_run()
        rn.text = "N"; rn.font.size = Pt(11); rn.font.bold = True
        rn.font.color.rgb = RGBColor(40, 40, 40); rn.font.name = "맑은 고딕"
        ntx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_scalebar(self, slide, sw_cm, sh_cm, r_km_y):
        """좌하단 스케일바 (2km)."""
        sb_km = 2.0
        sb_len_cm = (sb_km / (r_km_y * 2)) * sh_cm
        sb_x = Cm(0.6)
        sb_y = Cm(sh_cm) - Cm(0.9)
        bar_h = Cm(0.18)

        # 흰 반투명 배경
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            sb_x - Cm(0.15), sb_y - Cm(0.1),
            Cm(sb_len_cm) + Cm(0.3), Cm(0.8),
        )
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg.fill.transparency = 0.15
        bg.line.color.rgb = RGBColor(180, 180, 180); bg.line.width = Pt(0.5)

        # 막대 (흑백 2단)
        sb_half = Cm(sb_len_cm / 2)
        b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sb_x, sb_y, sb_half, bar_h)
        b1.fill.solid(); b1.fill.fore_color.rgb = RGBColor(40, 40, 40)
        b1.line.color.rgb = RGBColor(40, 40, 40); b1.line.width = Pt(0.4)
        b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sb_x + sb_half, sb_y, sb_half, bar_h)
        b2.fill.solid(); b2.fill.fore_color.rgb = RGBColor(255, 255, 255)
        b2.line.color.rgb = RGBColor(40, 40, 40); b2.line.width = Pt(0.4)

        # 라벨
        for ratio, label in [(0.0, "0"), (0.5, "1"), (1.0, "2 km")]:
            lx = sb_x + Cm(sb_len_cm) * ratio - Cm(0.5)
            tx = slide.shapes.add_textbox(lx, sb_y + bar_h + Cm(0.02),
                                          Cm(1.0), Cm(0.45))
            self._zero_margins(tx)
            r = tx.text_frame.paragraphs[0].add_run()
            r.text = label; r.font.size = Pt(8); r.font.bold = True
            r.font.color.rgb = RGBColor(40, 40, 40); r.font.name = "맑은 고딕"
            tx.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_pois(self, slide, query, drawn, bbox_str,
                  eMinX, eMaxX, eMinY, eMaxY, SW, SH, icon, color):
        """POI 라벨 추가 (IC/JC/역)."""
        pois = self._search_pois(bbox_str, query, size=50)
        col_rgb = RGBColor(*color)
        for poi in pois:
            title = poi.get("title", "")
            if not title or title in drawn or len(title) > 12:
                continue
            try:
                lon = float(poi["point"]["x"])
                lat = float(poi["point"]["y"])
            except Exception:
                continue
            if not (eMinX < lon < eMaxX and eMinY < lat < eMaxY):
                continue

            px_inch = (lon - eMinX) / (eMaxX - eMinX) * SW
            py_inch = (eMaxY - lat) / (eMaxY - eMinY) * SH

            # 아이콘 (Inches)
            if icon == "ic_jc":
                tag = "IC" if ("IC" in title or "인터체인지" in title) else "JC"
                ic = slide.shapes.add_shape(
                    MSO_SHAPE.DIAMOND,
                    Inches(px_inch) - Inches(0.08), Inches(py_inch) - Inches(0.08),
                    Inches(0.16), Inches(0.16),
                )
                ic.fill.solid(); ic.fill.fore_color.rgb = col_rgb
                ic.line.color.rgb = RGBColor(255, 255, 255); ic.line.width = Pt(0.8)
                # IC/JC 글자 (내부)
                self._zero_margins(ic)
                ri = ic.text_frame.paragraphs[0].add_run()
                ri.text = tag
                ri.font.size = Pt(5); ri.font.bold = True
                ri.font.color.rgb = RGBColor(255, 255, 255); ri.font.name = "맑은 고딕"
                ic.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            elif icon == "station":
                ic = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(px_inch) - Inches(0.04), Inches(py_inch) - Inches(0.04),
                    Inches(0.08), Inches(0.08),
                )
                ic.fill.solid(); ic.fill.fore_color.rgb = RGBColor(255, 255, 255)
                ic.line.color.rgb = RGBColor(40, 40, 40); ic.line.width = Pt(0.8)

            # 텍스트 라벨 (아이콘 우측)
            box_w = Cm(0.5 * len(title) + 0.3)
            box_h = Cm(0.5)
            tx = slide.shapes.add_textbox(
                Inches(px_inch) + Inches(0.08), Inches(py_inch) - Inches(0.05),
                box_w, box_h,
            )
            self._zero_margins(tx)
            tx.text_frame.word_wrap = False
            tx.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            rn = tx.text_frame.paragraphs[0].add_run()
            # 불필요 접두 제거
            clean_title = title.replace("인터체인지", "").replace("지하철 ", "").strip()
            rn.text = clean_title or title
            rn.font.size = Pt(7); rn.font.bold = True
            rn.font.color.rgb = col_rgb; rn.font.name = "맑은 고딕"
            self._text_outline_white(rn)
            drawn.add(title)
