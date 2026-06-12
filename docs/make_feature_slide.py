# -*- coding: utf-8 -*-
"""KH LandHub 기능 설명 1페이지 PPTX 생성 (스크린샷 + 콜아웃 + 지시선)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "landhub_screen.png")
LOGO = os.path.join(BASE, "..", "scratch", "deck_assets", "logo_kunhwa.png")
OUT = os.path.join(BASE, "KH_LandHub_기능소개.pptx")

NAVY = RGBColor(0x1F, 0x38, 0x64)
NAVY_L = RGBColor(0xEE, 0xF2, 0xF8)
ORANGE = RGBColor(0xC0, 0x56, 0x1E)
ORANGE_L = RGBColor(0xFB, 0xF1, 0xE7)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
SW, SH = 13.333, 7.5

def rect(x, y, w, h, fill, line=None, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def textbox(x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, wrap=True):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold, color, font)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2); p.space_before = Pt(0)
        p.line_spacing = 1.05
        for (text, size, bold, color, font) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb

def connector(x1, y1, x2, y2, color, width=1.5):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    # 지시선 끝(이미지 쪽)에 화살표 추가
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    cn.shadow.inherit = False
    return cn

def dot(xc, yc, color, d=0.12):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xc-d/2), Inches(yc-d/2), Inches(d), Inches(d))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.color.rgb = WHITE; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

FONT = "맑은 고딕"

# ---- 배경 ----
rect(0, 0, SW, SH, WHITE)
# 상단 네이비 밴드
rect(0, 0, SW, 1.5, NAVY, shape=MSO_SHAPE.RECTANGLE)
# 하단 얇은 네이비 라인 밴드 (이미지 아래 태그라인 바)

# ---- 제목 ----
textbox(0.55, 0.18, 9.8, 0.75, [
    [("KH LandHub ", 30, True, WHITE, FONT), ("화면 구성 및 주요 기능", 30, True, RGBColor(0xCA,0xDC,0xFC), FONT)],
], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
textbox(0.57, 0.92, 10.5, 0.45, [
    [("구역계 한 장으로 현황분석부터 인허가 도서 작성까지 — 하나의 화면에서 처리", 13.5, False, RGBColor(0xD6,0xDE,0xEE), FONT)],
], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
# 로고 (우상단)
if os.path.exists(LOGO):
    slide.shapes.add_picture(LOGO, Inches(12.05), Inches(0.42), height=Inches(0.62))

# ---- 스크린샷 ----
IMG_X, IMG_Y, IMG_W = 2.77, 1.85, 7.79
IMG_H = IMG_W * 863.0 / 1798.0  # = 3.738
# 이미지 테두리 프레임
rect(IMG_X-0.04, IMG_Y-0.04, IMG_W+0.08, IMG_H+0.08, WHITE, line=NAVY, line_w=1.25, shape=MSO_SHAPE.RECTANGLE)
slide.shapes.add_picture(IMG, Inches(IMG_X), Inches(IMG_Y), Inches(IMG_W), Inches(IMG_H))

def fx(frac):  # 이미지 가로 비율 → 슬라이드 X
    return IMG_X + frac * IMG_W
def fy(frac):  # 이미지 세로 비율 → 슬라이드 Y
    return IMG_Y + frac * IMG_H

# ---- 좌측 콜아웃 (사이드바) : 네이비 ----
BOX_W = 2.42
LX = 0.22
left_items = [
    ("🔍  장소 검색", "장소명을 입력하면 해당 위치로 지도가 즉시 이동", fy(0.10), fx(0.02)),
    ("①  구역계 범위 지정", "DXF 업로드 또는 지도에서 직접 그리기로 사업지구 경계 지정 · 국내 전 좌표계 자동 변환", fy(0.20), fx(0.02)),
    ("②  지도 레이어 · 데이터 추출", "브이월드·국립생태원 등 공공 공간정보를 표시하고 DXF/SHP로 일괄 추출", fy(0.46), fx(0.02)),
    ("③  현황 분석 보고서", "토지대장 자동분석(Excel), 공적장부(산지조서·협의서·소유자구분도), 현황분석 보고서(HWP) 자동 생성", fy(0.75), fx(0.02)),
    ("④  기타 도구", "SHP → DXF 변환 등 부가 기능 제공", fy(0.97), fx(0.02)),
]
# 박스 세로 균등 배치
box_h = 0.66
top0 = 1.86
gap = (5.60 - top0 - box_h) / (len(left_items) - 1)  # 마지막 박스 하단 5.60
for i, (head, body, ty, tx) in enumerate(left_items):
    by = top0 + i * gap
    cy = by + box_h/2
    rect(LX, by, BOX_W, box_h, NAVY_L, line=NAVY, line_w=0.75, radius=0.10)
    textbox(LX, by, BOX_W, box_h, [
        [(head, 10.5, True, NAVY, FONT)],
        [(body, 8.3, False, DARK, FONT)],
    ], anchor=MSO_ANCHOR.MIDDLE)
    connector(LX+BOX_W, cy, tx, ty, NAVY, 1.4)
    dot(tx, ty, NAVY)

# ---- 우측 콜아웃 (지도/결과) : 오렌지 ----
RX = SW - 0.22 - BOX_W
right_items = [
    ("📋  분석 결과 패널", "분석 진행 상황과 결과를 실시간으로 안내", 2.30, 0.85, fx(0.86), fy(0.22)),
    ("🗺️  인터랙티브 지도", "일반·위성·하이브리드 배경 위에 분석 레이어를 겹쳐 시각화", 3.95, 0.95, fx(0.70), fy(0.62)),
]
for (head, body, by, bh, tx, ty) in right_items:
    cy = by + bh/2
    rect(RX, by, BOX_W, bh, ORANGE_L, line=ORANGE, line_w=0.75, radius=0.10)
    textbox(RX, by, BOX_W, bh, [
        [(head, 10.5, True, ORANGE, FONT)],
        [(body, 8.3, False, DARK, FONT)],
    ], anchor=MSO_ANCHOR.MIDDLE)
    connector(RX, cy, tx, ty, ORANGE, 1.4)
    dot(tx, ty, ORANGE)

# ---- 하단 태그라인 바 ----
rect(IMG_X, 6.06, IMG_W, 0.66, NAVY, radius=0.18)
textbox(IMG_X, 6.06, IMG_W, 0.66, [
    [("전문가가 아니어도, 구역계 한 장으로 사업지구의 토지·규제·환경 여건을 자동 분석", 12.5, True, WHITE, FONT)],
], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---- 입력→출력 미니 플로우 (하단 좌측 여백) ----
textbox(0.22, 6.06, 2.42, 0.66, [
    [("입력", 9, True, GRAY, FONT), ("  구역계 DXF 1장", 9, False, DARK, FONT)],
    [("출력", 9, True, GRAY, FONT), ("  조서·보고서·도면", 9, False, DARK, FONT)],
], anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print("SAVED:", OUT)
